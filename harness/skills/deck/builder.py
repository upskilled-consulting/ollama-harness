"""Build a .pptx deck from DesignTheme + list of (name, markdown) pairs."""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .theme import DesignTheme

# ── Slide geometry ─────────────────────────────────────────────────────────────
W  = Inches(13.33)   # 16:9 widescreen width
H  = Inches(7.5)     # 16:9 widescreen height
MX = Inches(0.65)    # horizontal margin
MY = Inches(0.5)     # vertical margin
CW = W - 2 * MX     # content width

MAX_BULLETS = 6


# ── Low-level helpers ──────────────────────────────────────────────────────────
def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = h[0] * 2 + h[1] * 2 + h[2] * 2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _blank_slide(prs):  # type: ignore[no-untyped-def]
    try:
        layout = prs.slide_layouts[6]
    except IndexError:
        layout = prs.slide_layouts[-1]
    return prs.slides.add_slide(layout)


def _textbox(
    slide, text: str,
    left, top, width, height,
    font_name: str, size: int, color: str,
    bold: bool = False, italic: bool = False,
    align=PP_ALIGN.LEFT,
):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text         = text
    run.font.name    = font_name
    run.font.size    = Pt(size)
    run.font.color.rgb = _rgb(color)
    run.font.bold    = bold
    run.font.italic  = italic
    return txb


def _rect(slide, left, top, width, height, fill_color: str):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.fill.background()
    return shape


def _bullets_box(slide, bullets: list[str], left, top, width, height, theme: DesignTheme, size: int = 14):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text          = f"•  {text}"
        run.font.name     = theme.body_font
        run.font.size     = Pt(size)
        run.font.color.rgb = _rgb(theme.text_primary)
    return txb


# ── Slide spec ─────────────────────────────────────────────────────────────────
@dataclass
class SlideSpec:
    kind: Literal["title", "section", "bullets", "callout", "table"]
    title:         str             = ""
    subtitle:      str             = ""
    bullets:       list[str]       = field(default_factory=list)
    callout:       str             = ""
    table_headers: list[str]       = field(default_factory=list)
    table_rows:    list[list[str]] = field(default_factory=list)
    source_name:   str             = ""


# ── Markdown → SlideSpec list ──────────────────────────────────────────────────
def _flush_bullets(specs: list[SlideSpec], title: str, bullets: list[str], source: str) -> None:
    if not bullets:
        return
    for i in range(0, max(len(bullets), 1), MAX_BULLETS):
        chunk  = bullets[i:i + MAX_BULLETS]
        suffix = " (cont.)" if i > 0 else ""
        specs.append(SlideSpec(kind="bullets", title=title + suffix, bullets=chunk, source_name=source))


def _parse_table_lines(lines: list[str]) -> tuple[list[str], list[list[str]]]:
    headers: list[str]       = []
    rows:    list[list[str]] = []
    for line in lines:
        if re.match(r"^[-:| ]+$", line.strip()):
            continue
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if not headers:
            headers = cells
        else:
            rows.append(cells)
    return headers, rows


def _clean_inline(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*",     r"\1", text)
    text = re.sub(r"`(.+?)`",       r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    return text.strip()


def markdown_to_slides(name: str, md: str) -> list[SlideSpec]:
    specs:           list[SlideSpec] = []
    current_title:   str             = name
    current_bullets: list[str]       = []
    table_lines:     list[str]       = []
    in_table         = False
    paras_this_section = 0   # paragraphs auto-promoted to bullets in current section

    lines = md.splitlines()
    i = 0
    while i < len(lines):
        line     = lines[i]
        stripped = line.strip()

        # ── Table ─────────────────────────────────────────────────────────────
        if stripped.startswith("|"):
            if not in_table:
                _flush_bullets(specs, current_title, current_bullets, name)
                current_bullets = []
                in_table    = True
                table_lines = [stripped]
            else:
                table_lines.append(stripped)
            i += 1
            continue
        if in_table:
            headers, rows = _parse_table_lines(table_lines)
            if headers and rows:
                specs.append(SlideSpec(
                    kind="table", title=current_title,
                    table_headers=headers, table_rows=rows[:10],
                    source_name=name,
                ))
            in_table    = False
            table_lines = []
            continue   # re-process current line

        # ── H1 → title slide ──────────────────────────────────────────────────
        if re.match(r"^# [^#]", stripped):
            _flush_bullets(specs, current_title, current_bullets, name)
            current_bullets    = []
            paras_this_section = 0
            current_title      = stripped[2:].strip()
            subtitle = ""
            j = i + 1
            while j < len(lines) and not lines[j].strip():
                j += 1
            if j < len(lines) and not lines[j].strip().startswith("#"):
                subtitle = _clean_inline(lines[j].strip())
                i = j
            specs.append(SlideSpec(kind="title", title=current_title, subtitle=subtitle, source_name=name))
            i += 1
            continue

        # ── H2 → content section ──────────────────────────────────────────────
        if stripped.startswith("## "):
            _flush_bullets(specs, current_title, current_bullets, name)
            current_bullets    = []
            paras_this_section = 0
            current_title      = stripped[3:].strip()
            # Peek: if immediately followed by another heading → standalone section divider
            j = i + 1
            while j < len(lines) and not lines[j].strip():
                j += 1
            if j >= len(lines) or lines[j].strip().startswith("#"):
                specs.append(SlideSpec(kind="section", title=current_title, source_name=name))
            i += 1
            continue

        # ── H3 → sub-section label ────────────────────────────────────────────
        if stripped.startswith("### "):
            _flush_bullets(specs, current_title, current_bullets, name)
            current_bullets    = []
            paras_this_section = 0
            current_title      = stripped[4:].strip()
            i += 1
            continue

        # ── H4+ → bold bullet ─────────────────────────────────────────────────
        if stripped.startswith("#### "):
            current_bullets.append(stripped[5:].strip())
            i += 1
            continue

        # ── Blockquote → callout ──────────────────────────────────────────────
        if stripped.startswith("> "):
            _flush_bullets(specs, current_title, current_bullets, name)
            current_bullets = []
            specs.append(SlideSpec(
                kind="callout", title=current_title,
                callout=_clean_inline(stripped[2:]),
                source_name=name,
            ))
            i += 1
            continue

        # ── List item ─────────────────────────────────────────────────────────
        if re.match(r"^[-*+]\s+", stripped) or re.match(r"^\d+\.\s+", stripped):
            text = re.sub(r"^[-*+]\s+|\d+\.\s+", "", stripped)
            text = _clean_inline(text)
            if text:
                current_bullets.append(text[:180] + ("…" if len(text) > 180 else ""))
            i += 1
            continue

        # ── Paragraph → bullet (max 2 per section, min 80 chars) ────────────
        if (stripped and len(stripped) > 80
                and not stripped.startswith("```")
                and current_title
                and paras_this_section < 2):
            text = _clean_inline(stripped)
            current_bullets.append(text[:180] + ("..." if len(text) > 180 else ""))
            paras_this_section += 1

        i += 1

    # Flush any table still open at EOF
    if in_table and table_lines:
        headers, rows = _parse_table_lines(table_lines)
        if headers and rows:
            specs.append(SlideSpec(
                kind="table", title=current_title,
                table_headers=headers, table_rows=rows[:10],
                source_name=name,
            ))

    _flush_bullets(specs, current_title, current_bullets, name)
    return specs


# ── Slide renderers ────────────────────────────────────────────────────────────
def _render_title(slide, spec: SlideSpec, t: DesignTheme) -> None:
    _set_bg(slide, t.bg_primary)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.07), t.accent)
    _textbox(slide, spec.title,
             MX, Inches(2.2), CW, Inches(1.8),
             t.heading_font, t.title_size, t.accent, bold=True, align=PP_ALIGN.CENTER)
    if spec.subtitle:
        _textbox(slide, spec.subtitle,
                 MX, Inches(4.1), CW, Inches(0.9),
                 t.body_font, t.h3_size, t.text_secondary, align=PP_ALIGN.CENTER)
    if spec.source_name:
        _textbox(slide, spec.source_name,
                 W - Inches(3.6), H - Inches(0.38), Inches(3.3), Inches(0.3),
                 t.body_font, t.small_size, t.text_secondary, align=PP_ALIGN.RIGHT)


def _render_section(slide, spec: SlideSpec, t: DesignTheme) -> None:
    _set_bg(slide, t.bg_surface)
    _rect(slide, Inches(0), Inches(0), Inches(0.09), H, t.accent)
    _textbox(slide, spec.title,
             Inches(0.55), Inches(2.7), W - Inches(0.75), Inches(1.6),
             t.heading_font, t.h1_size, t.accent, bold=True)
    if spec.subtitle:
        _textbox(slide, spec.subtitle,
                 Inches(0.55), Inches(4.5), W - Inches(0.75), Inches(0.7),
                 t.body_font, t.h3_size, t.text_secondary)


def _render_bullets(slide, spec: SlideSpec, t: DesignTheme) -> None:
    _set_bg(slide, t.bg_primary)
    _textbox(slide, spec.title,
             MX, MY, CW, Inches(0.65),
             t.heading_font, t.h2_size, t.text_primary, bold=True)
    # Thin accent rule below title
    _rect(slide, MX, MY + Inches(0.68), CW, Emu(38000), t.accent)
    _bullets_box(slide, spec.bullets,
                 MX, MY + Inches(0.82), CW, H - MY - Inches(1.0),
                 t, size=t.body_size)


def _render_callout(slide, spec: SlideSpec, t: DesignTheme) -> None:
    _set_bg(slide, t.bg_surface)
    _rect(slide, Inches(0), Inches(0), W, Inches(0.06), t.accent)
    text = spec.callout or spec.title
    _textbox(slide, f'"{text}"',
             MX, Inches(1.9), CW, Inches(3.8),
             t.heading_font, 22, t.accent, italic=True, align=PP_ALIGN.CENTER)
    if spec.title and spec.callout:
        _textbox(slide, f"— {spec.title}",
                 MX, Inches(5.8), CW, Inches(0.5),
                 t.body_font, t.small_size, t.text_secondary, align=PP_ALIGN.CENTER)


def _render_table(slide, spec: SlideSpec, t: DesignTheme) -> None:
    _set_bg(slide, t.bg_primary)
    _textbox(slide, spec.title,
             MX, MY, CW, Inches(0.6),
             t.heading_font, t.h2_size, t.text_primary, bold=True)

    headers = spec.table_headers
    rows    = spec.table_rows[:8]
    if not headers:
        return

    n_cols = len(headers)
    n_rows = len(rows) + 1
    table  = slide.shapes.add_table(
        n_rows, n_cols,
        MX, MY + Inches(0.75), CW, Inches(5.5),
    ).table

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(t.bg_surface)
        p = cell.text_frame.paragraphs[0]
        for run in p.runs:
            run.font.name  = t.heading_font
            run.font.size  = Pt(11)
            run.font.bold  = True
            run.font.color.rgb = _rgb(t.accent)

    # Body rows — alternating background
    for ri, row in enumerate(rows):
        row_color = t.bg_primary if ri % 2 == 0 else t.bg_surface
        for ci in range(n_cols):
            cell = table.cell(ri + 1, ci)
            cell.text = row[ci] if ci < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(row_color)
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.name  = t.body_font
                run.font.size  = Pt(10)
                run.font.color.rgb = _rgb(t.text_primary)


_RENDERERS = {
    "title":   _render_title,
    "section": _render_section,
    "bullets": _render_bullets,
    "callout": _render_callout,
    "table":   _render_table,
}


# ── Entry point ────────────────────────────────────────────────────────────────
def build_deck(
    theme: DesignTheme,
    content: list[tuple[str, str]],
    output_path: str,
    deck_title: str = "",
) -> str:
    """Build and save a .pptx. Returns the resolved output path."""
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    all_specs: list[SlideSpec] = []

    if deck_title:
        all_specs.append(SlideSpec(kind="title", title=deck_title))

    for name, md in content:
        specs = markdown_to_slides(name, md)
        # Insert a section divider between documents in multi-file decks
        if len(content) > 1 and all_specs and specs and specs[0].kind not in ("title", "section"):
            all_specs.append(SlideSpec(kind="section", title=name))
        all_specs.extend(specs)

    print(f"  [deck] rendering {len(all_specs)} slides...")

    for spec in all_specs:
        s = _blank_slide(prs)
        _RENDERERS.get(spec.kind, _render_bullets)(s, spec, theme)

    out = Path(output_path).expanduser()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"  [deck] saved -> {out}")
    return str(out)
